import os
import io
import logging
from typing import Dict, Any, Optional, List, Tuple
import numpy as np
from PIL import Image
import cv2
import pytesseract
import fitz
import pdf2image
from skimage import filters, morphology
import chardet
from settings import settings

logger = logging.getLogger(__name__)


class OCRProcessor:
    """
    A synchronous processor for OCR tasks. All its methods run directly.
    Designed to be used inside a dedicated worker process.
    """

    def __init__(self, config: Optional[Dict] = None):
        self.config = config or {}
        self.use_tesseract = self.config.get("use_tesseract", True)
        self.use_easyocr = self.config.get("use_easyocr", True)
        self.use_keras_ocr = self.config.get("use_keras_ocr", False)
        self._ocr_engines: Dict[str, Any] = {}

    @property
    def easyocr_engine(self):
        """Lazy loader for the EasyOCR engine."""
        if "easyocr" not in self._ocr_engines:
            if self.use_easyocr:
                try:
                    logger.info("Initializing EasyOCR. This may take a while...")
                    import easyocr

                    # Note: We set torch threads in the worker, not here.
                    self._ocr_engines["easyocr"] = easyocr.Reader(["en"], gpu=False)
                    logger.info("EasyOCR initialized successfully.")
                except Exception as e:
                    logger.error(f"Failed to initialize EasyOCR: {e}")
                    self._ocr_engines["easyocr"] = None
            else:
                self._ocr_engines["easyocr"] = None
        return self._ocr_engines.get("easyocr")

    @property
    def keras_ocr_engine(self):
        """Lazy loader for the Keras-OCR engine."""
        if "keras_ocr" not in self._ocr_engines:
            if self.use_keras_ocr:
                try:
                    logger.info("Initializing Keras-OCR. This may take a while...")
                    import keras_ocr

                    self._ocr_engines["keras_ocr"] = keras_ocr.pipeline.Pipeline()
                    logger.info("Keras-OCR initialized successfully.")
                except Exception as e:
                    logger.error(f"Failed to initialize Keras-OCR: {e}")
                    self._ocr_engines["keras_ocr"] = None
            else:
                self._ocr_engines["keras_ocr"] = None
        return self._ocr_engines.get("keras_ocr")

    def preprocess_image(self, image: np.ndarray) -> Tuple[np.ndarray, Dict]:
        """Apply image preprocessing to improve OCR accuracy."""
        preprocessing_info = {}
        try:
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image.copy()
            denoised = cv2.fastNlMeansDenoising(gray)
            thresh = cv2.adaptiveThreshold(
                denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
            )
            kernel = np.ones((1, 1), np.uint8)
            cleaned = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
            corrected = self._correct_skew(cleaned)
            preprocessing_info = {
                "applied_denoising": True,
                "applied_thresholding": True,
                "applied_morphology": True,
                "applied_skew_correction": True,
            }
            return corrected, preprocessing_info
        except Exception as e:
            logger.warning(f"Image preprocessing failed: {e}")
            return image, {"error": str(e)}

    def _correct_skew(self, image: np.ndarray) -> np.ndarray:
        """Correct skew in the image."""
        try:
            edges = cv2.Canny(image, 50, 150, apertureSize=3)
            lines = cv2.HoughLinesP(
                edges, 1, np.pi / 180, threshold=100, minLineLength=100, maxLineGap=10
            )
            if lines is not None:
                angles = [
                    np.arctan2(y2 - y1, x2 - x1)
                    for line in lines
                    for x1, y1, x2, y2 in line
                ]
                median_angle = np.median(angles)
                angle_deg = np.degrees(median_angle)
                if abs(angle_deg) > 0.5:
                    (h, w) = image.shape[:2]
                    center = (w // 2, h // 2)
                    M = cv2.getRotationMatrix2D(center, angle_deg, 1.0)
                    corrected = cv2.warpAffine(
                        image,
                        M,
                        (w, h),
                        flags=cv2.INTER_CUBIC,
                        borderMode=cv2.BORDER_REPLICATE,
                    )
                    return corrected
            return image
        except Exception as e:
            logger.warning(f"Skew correction failed: {e}")
            return image

    def extract_text_tesseract(
        self, image: np.ndarray, config: str = "--psm 6"
    ) -> Dict:
        """Extract text using Tesseract OCR."""
        try:
            data = pytesseract.image_to_data(
                image, config=config, output_type=pytesseract.Output.DICT
            )
            min_confidence = self.config.get("min_confidence", 30)
            confident_text = []
            for i, conf in enumerate(data["conf"]):
                if int(conf) > min_confidence:
                    text = data["text"][i].strip()
                    if text:
                        confident_text.append(
                            {
                                "text": text,
                                "confidence": int(conf),
                                "bbox": (
                                    data["left"][i],
                                    data["top"][i],
                                    data["width"][i],
                                    data["height"][i],
                                ),
                            }
                        )
            full_text = " ".join([item["text"] for item in confident_text])
            return {
                "text": full_text,
                "words": confident_text,
                "engine": "tesseract",
                "success": True,
            }
        except Exception as e:
            logger.error(f"Tesseract OCR failed: {e}")
            return {
                "text": "",
                "words": [],
                "engine": "tesseract",
                "success": False,
                "error": str(e),
            }

    def extract_text_easyocr(self, image: np.ndarray) -> Dict:
        """Extract text using EasyOCR."""
        engine = self.easyocr_engine
        try:
            if not engine:
                return {
                    "text": "",
                    "words": [],
                    "engine": "easyocr",
                    "success": False,
                    "error": "EasyOCR not initialized.",
                }
            results = engine.readtext(image)
            words = []
            full_text_parts = []
            for bbox, text, confidence in results:
                if confidence > self.config.get("min_confidence", 0.3):
                    words.append(
                        {"text": text, "confidence": float(confidence), "bbox": bbox}
                    )
                    full_text_parts.append(text)
            return {
                "text": " ".join(full_text_parts),
                "words": words,
                "engine": "easyocr",
                "success": True,
            }
        except Exception as e:
            logger.error(f"EasyOCR failed: {e}")
            return {
                "text": "",
                "words": [],
                "engine": "easyocr",
                "success": False,
                "error": str(e),
            }

    def extract_text_keras_ocr(self, image: np.ndarray) -> Dict:
        """Extract text using Keras-OCR."""
        engine = self.keras_ocr_engine
        try:
            if not engine:
                return {
                    "text": "",
                    "words": [],
                    "engine": "keras_ocr",
                    "success": False,
                    "error": "Keras-OCR not initialized.",
                }
            if len(image.shape) == 2:  # Convert grayscale to RGB
                image = cv2.cvtColor(image, cv2.COLOR_GRAY2RGB)
            prediction_groups = engine.recognize([image])
            words = [
                {"text": text, "confidence": 1.0, "bbox": bbox.tolist()}
                for text, bbox in prediction_groups[0]
            ]
            full_text = " ".join([word["text"] for word in words])
            return {
                "text": full_text,
                "words": words,
                "engine": "keras_ocr",
                "success": True,
            }
        except Exception as e:
            logger.error(f"Keras-OCR failed: {e}")
            return {
                "text": "",
                "words": [],
                "engine": "keras_ocr",
                "success": False,
                "error": str(e),
            }

    def simple_layout_analysis(self, image: np.ndarray) -> Dict:
        """Simple layout analysis using OpenCV."""
        try:
            # Convert to grayscale if needed
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image.copy()

            # Find contours for text regions
            _, binary = cv2.threshold(
                gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU
            )
            contours, _ = cv2.findContours(
                binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
            )

            text_regions = []
            for contour in contours:
                x, y, w, h = cv2.boundingRect(contour)
                if w > 10 and h > 10:  # Filter small noise
                    text_regions.append({"x": x, "y": y, "width": w, "height": h})

            return {
                "text_regions": text_regions,
                "num_regions": len(text_regions),
                "success": True,
            }
        except Exception as e:
            logger.error(f"Layout analysis failed: {e}")
            return {"success": False, "error": str(e)}

    def ensemble_ocr(self, image: np.ndarray) -> Dict:
        """Run multiple OCR engines and combine results."""
        results = []
        if self.use_tesseract:
            tesseract_result = self.extract_text_tesseract(image)
            if tesseract_result["success"]:
                results.append(tesseract_result)
        if self.use_easyocr:
            easyocr_result = self.extract_text_easyocr(image)
            if easyocr_result["success"]:
                results.append(easyocr_result)
        if self.use_keras_ocr:
            keras_result = self.extract_text_keras_ocr(image)
            if keras_result["success"]:
                results.append(keras_result)
        if not results:
            return {"text": "", "success": False, "error": "All OCR engines failed"}
        best_result = max(results, key=lambda x: len(x["text"]))
        best_result["ensemble_results"] = [r["engine"] for r in results]
        return best_result

    def process_image(self, image_data: bytes, filename: str = "") -> Dict:
        """
        This is now the main synchronous method for processing a single image.
        It no longer spawns a subprocess.
        """
        try:
            image = Image.open(io.BytesIO(image_data))
            image_np = np.array(image)
            max_size = self.config.get("max_image_size", (2000, 2000))
            if image.size[0] > max_size[0] or image.size[1] > max_size[1]:
                image.thumbnail(max_size, Image.Resampling.LANCZOS)
                image_np = np.array(image)

            preprocessed_image, preprocess_info = self.preprocess_image(image_np)
            layout_info = self.simple_layout_analysis(preprocessed_image)
            ocr_result = self.ensemble_ocr(preprocessed_image)

            return {
                "success": ocr_result.get("success", False),
                "text": ocr_result.get("text", ""),
                "filename": filename,
                "preprocessing": preprocess_info,
                "layout_analysis": layout_info,
                "ocr_details": ocr_result,
                "image_size": image.size,
            }
        except Exception as e:
            logger.error(f"Image processing failed for {filename}: {e}", exc_info=True)
            return {"success": False, "text": "", "filename": filename, "error": str(e)}


class EnhancedDocumentProcessor:
    """
    This class orchestrates text extraction. It remains largely unchanged,
    but it now calls the synchronous `process_image` method.
    """

    def __init__(self):
        self.ocr_processor = OCRProcessor(
            {
                "use_tesseract": True,
                "use_easyocr": True,
                "use_keras_ocr": False,
                "min_confidence": 0.3,
                "max_image_size": (2000, 2000),
            }
        )
        self.supported_extensions = settings.SUPPORTED_EXTENSIONS

    def detect_encoding(self, file_content: bytes) -> str:
        """Detect file encoding using chardet."""
        try:
            result = chardet.detect(file_content)
            encoding = result.get("encoding", "utf-8")
            confidence = result.get("confidence", 0)

            # Fallback to utf-8 if confidence is too low
            if confidence < 0.7:
                encoding = "utf-8"

            return encoding
        except Exception as e:
            logger.warning(f"Encoding detection failed: {e}")
            return "utf-8"

    def extract_images_from_pdf(self, file_content: bytes, filename: str) -> List[Dict]:
        """Extract images from PDF and process them with OCR."""
        images_text = []
        try:
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            for page_num in range(min(len(pdf_document), 5)):
                page = pdf_document[page_num]
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_document, xref)
                        if pix.n - pix.alpha < 4:
                            img_data = pix.tobytes("png")
                            ocr_result = self.ocr_processor.process_image(
                                img_data,
                                f"{filename}_page_{page_num+1}_img_{img_index+1}.png",
                            )
                            if ocr_result["success"] and ocr_result["text"].strip():
                                images_text.append(
                                    {
                                        "page": page_num + 1,
                                        "image_index": img_index + 1,
                                        "text": ocr_result["text"],
                                        "ocr_details": ocr_result,
                                    }
                                )
                        pix = None
                    except Exception as e:
                        logger.warning(
                            f"Failed to process image {img_index} on page {page_num}: {e}"
                        )
            pdf_document.close()

            if not images_text:
                pages = pdf2image.convert_from_bytes(
                    file_content, dpi=200, first_page=1, last_page=3
                )
                for page_num, page_image in enumerate(pages):
                    img_bytes = io.BytesIO()
                    page_image.save(img_bytes, format="PNG")
                    img_data = img_bytes.getvalue()
                    ocr_result = self.ocr_processor.process_image(
                        img_data, f"{filename}_page_{page_num+1}_full.png"
                    )
                    if ocr_result["success"] and ocr_result["text"].strip():
                        images_text.append(
                            {
                                "page": page_num + 1,
                                "image_index": "full_page",
                                "text": ocr_result["text"],
                                "ocr_details": ocr_result,
                            }
                        )
        except Exception as e:
            logger.error(f"Failed to extract images from PDF {filename}: {e}")
        return images_text

    def extract_text_from_pdf_with_ocr(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from PDF using OCR on images."""
        try:
            # First try extracting text directly from PDF
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            text_content = ""

            for page in pdf_document:
                page_text = page.get_text()
                text_content += page_text + "\n"

            pdf_document.close()

            # If we got substantial text, return it
            if len(text_content.strip()) > 100:
                return text_content

            # Otherwise, fall back to OCR
            logger.info(f"PDF {filename} has minimal text, using OCR fallback")
            images_text = self.extract_images_from_pdf(file_content, filename)

            if images_text:
                ocr_text = "\n\n".join([img["text"] for img in images_text])
                return text_content + "\n\n--- OCR Content ---\n\n" + ocr_text

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from PDF {filename}: {e}")
            return None

    def extract_text_from_image_file(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from image files using OCR."""
        try:
            ocr_result = self.ocr_processor.process_image(file_content, filename)
            if ocr_result["success"]:
                return ocr_result["text"]
            return None
        except Exception as e:
            logger.error(f"Failed to extract text from image {filename}: {e}")
            return None

    def extract_text_from_file(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from various file types."""
        try:
            _, ext = os.path.splitext(filename.lower())

            # PDF files
            if ext == ".pdf":
                return self.extract_text_from_pdf_with_ocr(file_content, filename)

            # Image files
            elif ext in [
                ".png",
                ".jpg",
                ".jpeg",
                ".tiff",
                ".tif",
                ".bmp",
                ".gif",
                ".webp",
            ]:
                return self.extract_text_from_image_file(file_content, filename)

            # Document files
            elif ext in [".docx", ".doc"]:
                return self._extract_text_from_docx(file_content, filename)
            elif ext in [".pptx", ".ppt"]:
                return self._extract_text_from_pptx(file_content, filename)
            elif ext in [".xlsx", ".xls"]:
                return self._extract_text_from_excel(file_content, filename)
            elif ext in [".csv", ".tsv"]:
                return self._extract_text_from_csv(file_content, filename)
            elif ext in [".html", ".htm"]:
                return self._extract_text_from_html(file_content, filename)
            elif ext == ".eml":
                return self._extract_text_from_email(file_content, filename)

            # Plain text files
            else:
                return self._extract_text_from_plain_text(file_content, filename)

        except Exception as e:
            logger.error(f"Failed to extract text from {filename}: {e}")
            return None

    def process_file_content(
        self, file_content: bytes, filename: str, max_size_mb: int = 10
    ) -> Dict[str, Any]:
        """Process file content directly without mock objects."""
        result = {
            "success": False,
            "filename": filename,
            "text_content": None,
            "file_size": len(file_content),
            "file_type": None,
            "file_type_description": None,
            "error_message": None,
        }

        try:
            # Get file extension and description
            _, ext = os.path.splitext(filename.lower())
            result["file_type"] = ext
            result["file_type_description"] = self.get_file_type_description(filename)

            # Validate file size
            if not self.validate_file_size(result["file_size"], max_size_mb):
                result["error_message"] = (
                    f"File size ({result['file_size']/1024/1024:.1f} MB) exceeds maximum allowed size ({max_size_mb} MB)"
                )
                return result

            # Check if file type is supported
            if not self.is_supported_file(filename):
                result["error_message"] = (
                    f"File type '{ext}' is not supported. Supported types: {', '.join(sorted(self.supported_extensions))}"
                )
                return result

            # Extract text content
            text_content = self.extract_text_from_file(file_content, filename)

            if text_content is None:
                result["error_message"] = "Failed to extract text content from file"
                return result

            if len(text_content.strip()) == 0:
                result["error_message"] = (
                    "File appears to be empty or contains no extractable text"
                )
                return result

            result["text_content"] = text_content
            result["success"] = True

            logger.info(
                f"Successfully processed {result['file_type_description']}: {filename} ({result['file_size']} bytes)"
            )

        except Exception as e:
            result["error_message"] = f"Error processing file: {str(e)}"
            logger.error(f"Error processing file {filename}: {e}")

        return result

    def is_supported_file(self, filename: str) -> bool:
        """Check if file type is supported (including images)."""
        _, ext = os.path.splitext(filename.lower())
        return ext in self.supported_extensions

    def get_file_type_description(self, filename: str) -> str:
        """Get human-readable file type description including images."""
        _, ext = os.path.splitext(filename.lower())

        type_descriptions = {
            # Documents
            ".pdf": "PDF Document",
            ".docx": "Word Document",
            ".doc": "Word Document (Legacy)",
            ".pptx": "PowerPoint Presentation",
            ".ppt": "PowerPoint Presentation (Legacy)",
            ".xlsx": "Excel Spreadsheet",
            ".xls": "Excel Spreadsheet (Legacy)",
            ".csv": "CSV Data File",
            ".tsv": "Tab-Separated Values",
            ".html": "HTML Document",
            ".htm": "HTML Document",
            ".eml": "Email Message",
            ".txt": "Text File",
            ".md": "Markdown Document",
            ".json": "JSON Data File",
            ".xml": "XML Document",
            ".yaml": "YAML Configuration",
            ".yml": "YAML Configuration",
            # Images
            ".png": "PNG Image",
            ".jpg": "JPEG Image",
            ".jpeg": "JPEG Image",
            ".tiff": "TIFF Image",
            ".tif": "TIFF Image",
            ".bmp": "Bitmap Image",
            ".gif": "GIF Image",
            ".webp": "WebP Image",
        }

        return type_descriptions.get(
            ext, f"{ext.upper()} File" if ext else "Unknown File Type"
        )

    def validate_file_size(self, file_size: int, max_size_mb: int = 10) -> bool:
        """Validate file size."""
        max_size_bytes = max_size_mb * 1024 * 1024
        return file_size <= max_size_bytes

    # Additional helper methods for document types
    def _extract_text_from_docx(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from Word documents."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(file_content))
            text_content = ""

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content += " | ".join(row_text) + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from Word document {filename}: {e}")
            return None

    def _extract_text_from_pptx(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from PowerPoint presentations."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_content))
            text_content = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n\n--- Slide {slide_num} ---\n\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from PowerPoint {filename}: {e}")
            return None

    def _extract_text_from_excel(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from Excel files."""
        try:
            import openpyxl
            import xlrd

            # Try with openpyxl first (for .xlsx)
            if filename.lower().endswith(".xlsx"):
                workbook = openpyxl.load_workbook(
                    io.BytesIO(file_content), read_only=True
                )
                text_content = ""

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content += f"\n\n--- Sheet: {sheet_name} ---\n\n"

                    rows_data = []
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [
                            str(cell) if cell is not None else "" for cell in row
                        ]
                        if any(cell.strip() for cell in row_data):
                            rows_data.append(" | ".join(row_data))

                    text_content += "\n".join(rows_data)

                return text_content if text_content.strip() else None

            else:
                # Use xlrd for .xls files
                workbook = xlrd.open_workbook(file_contents=file_content)
                text_content = ""

                for sheet in workbook.sheets():
                    text_content += f"\n\n--- Sheet: {sheet.name} ---\n\n"

                    for row_idx in range(sheet.nrows):
                        row_data = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            row_data.append(str(cell_value) if cell_value else "")

                        if any(cell.strip() for cell in row_data):
                            text_content += " | ".join(row_data) + "\n"

                return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from Excel file {filename}: {e}")
            return None

    def _extract_text_from_csv(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from CSV files."""
        try:
            import pandas as pd

            # Detect encoding
            encoding = self.detect_encoding(file_content)

            # Try different separators
            separators = [",", ";", "\t", "|"]

            for sep in separators:
                try:
                    df = pd.read_csv(
                        io.BytesIO(file_content), encoding=encoding, sep=sep, nrows=5
                    )
                    if len(df.columns) > 1:
                        # Read the full file
                        df = pd.read_csv(
                            io.BytesIO(file_content), encoding=encoding, sep=sep
                        )

                        # Convert to text format
                        text_content = f"CSV File: {filename}\n\n"
                        text_content += (
                            f"Columns: {' | '.join(df.columns.tolist())}\n\n"
                        )

                        # Add sample of data
                        for idx, row in df.iterrows():
                            row_text = " | ".join(
                                [
                                    str(val) if pd.notna(val) else ""
                                    for val in row.values
                                ]
                            )
                            text_content += row_text + "\n"

                            # Limit to prevent huge files
                            if idx >= 1000:
                                text_content += f"\n... (truncated, showing first 1000 rows of {len(df)} total)\n"
                                break

                        return text_content

                except Exception:
                    continue

            # If all separators fail, treat as plain text
            return file_content.decode(encoding, errors="ignore")

        except Exception as e:
            logger.error(f"Failed to extract text from CSV {filename}: {e}")
            return None

    def _extract_text_from_html(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup

            encoding = self.detect_encoding(file_content)
            html_content = file_content.decode(encoding, errors="ignore")

            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text and clean it up
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text_content = " ".join(chunk for chunk in chunks if chunk)

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from HTML {filename}: {e}")
            return None

    def _extract_text_from_email(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from email files."""
        try:
            import email
            from email import policy
            from bs4 import BeautifulSoup

            # Parse email content
            msg = email.message_from_bytes(file_content, policy=policy.default)

            text_content = f"Email: {filename}\n\n"

            # Extract headers
            if msg["subject"]:
                text_content += f"Subject: {msg['subject']}\n"
            if msg["from"]:
                text_content += f"From: {msg['from']}\n"
            if msg["to"]:
                text_content += f"To: {msg['to']}\n"
            if msg["date"]:
                text_content += f"Date: {msg['date']}\n"

            text_content += "\n--- Email Body ---\n\n"

            # Extract body content
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_content()
                        if body:
                            text_content += body + "\n"
                    elif part.get_content_type() == "text/html":
                        # Fallback to HTML if no plain text
                        html_body = part.get_content()
                        if html_body and "text/plain" not in [
                            p.get_content_type() for p in msg.walk()
                        ]:
                            soup = BeautifulSoup(html_body, "html.parser")
                            text_content += soup.get_text() + "\n"
            else:
                # Single part message
                body = msg.get_content()
                if body:
                    text_content += body + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from email {filename}: {e}")
            return None

    def _extract_text_from_plain_text(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from plain text files with better encoding handling."""
        try:
            # Try to detect encoding first
            encoding = self.detect_encoding(file_content)

            # Decode with detected encoding
            text_content = file_content.decode(encoding, errors="ignore")

            # Clean up the text using ftfy for better unicode handling
            try:
                import ftfy

                text_content = ftfy.fix_text(text_content)
            except ImportError:
                # ftfy not available, continue without it
                pass

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from plain text file {filename}: {e}")
            return None
